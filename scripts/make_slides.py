#!/usr/bin/env python
"""Generate a fully-explanatory SERUM deck (~50 slides) with figures, worked
examples, and complete speaker notes on every slide."""
import os, struct
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = "/Users/amritha/serum"
FIG = os.path.join(ROOT, "paper", "figures")
RES = os.path.join(ROOT, "results")
OUT = os.path.join(ROOT, "presentation", "SERUM_slides.pptx")

NAVY  = RGBColor(0x1F, 0x35, 0x5E)
ACCENT= RGBColor(0x2E, 0x6F, 0xB0)
DARK  = RGBColor(0x25, 0x2B, 0x33)
GREY  = RGBColor(0x5A, 0x63, 0x6E)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
BOX   = RGBColor(0xEA, 0xF1, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1B, 0x7F, 0x4B)
AMBER = RGBColor(0xB5, 0x6A, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_n = [0]


def png_dim(p):
    with open(p, "rb") as f:
        f.read(16); return struct.unpack(">II", f.read(8))

def bg(s, c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def rect(s, l, t, w, h, color, line=None):
    sh = s.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(s, text, l, t, w, h, size=18, color=DARK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def para(tf, text, size=17, color=DARK, bold=False, italic=False, after=8, first=False, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(after)
    r = p.add_run(); r.text = ("•  " + text) if bullet else text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return p

def heading(s, title, kicker=None):
    bg(s, WHITE)
    if kicker:
        txt(s, kicker.upper(), Inches(0.6), Inches(0.40), Inches(12), Inches(0.35),
            size=12.5, color=ACCENT, bold=True)
    txt(s, title, Inches(0.6), Inches(0.70), Inches(12.15), Inches(0.95),
        size=27, color=NAVY, bold=True)
    rect(s, Inches(0), Inches(1.60), SW, Inches(0.06), ACCENT)
    _n[0] += 1
    txt(s, str(_n[0]), Inches(12.7), Inches(7.05), Inches(0.5), Inches(0.4),
        size=11, color=GREY, align=PP_ALIGN.RIGHT)

# ---------- slide constructors (all take notes) ----------
def divider(num, title, subtitle=""):
    s = prs.slides.add_slide(BLANK); bg(s, NAVY)
    txt(s, f"PART {num}", Inches(0.9), Inches(2.3), Inches(6), Inches(0.6),
        size=18, color=RGBColor(0x8F,0xA6,0xC8), bold=True)
    txt(s, title, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.6),
        size=40, color=WHITE, bold=True)
    if subtitle:
        txt(s, subtitle, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.2),
            size=19, color=RGBColor(0xC7,0xD6,0xEC), italic=True)
    notes(s, f"Section {num}: {title}. {subtitle}")
    return s

def explain(title, lead, paras, kicker=None, note="", lead_color=NAVY):
    s = prs.slides.add_slide(BLANK); heading(s, title, kicker)
    txt(s, lead, Inches(0.7), Inches(1.85), Inches(11.9), Inches(1.0),
        size=21, color=lead_color, bold=True)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.95), Inches(11.9), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(paras):
        if isinstance(item, tuple):
            para(tf, item[0], size=item[1] if len(item) > 1 else 17, first=(i == 0),
                 color=item[2] if len(item) > 2 else DARK, bullet=True)
        else:
            para(tf, item, size=17, first=(i == 0), bullet=True)
    notes(s, note); return s

def worked(title, intro, steps, conclusion, kicker="Worked example", note=""):
    s = prs.slides.add_slide(BLANK); heading(s, title, kicker)
    txt(s, intro, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.8), size=18,
        color=DARK)
    top = Inches(2.7)
    for i, (lab, body, col) in enumerate(steps):
        y = Emu(int(top) + i * int(Inches(0.92)))
        rect(s, Inches(0.7), y, Inches(0.9), Inches(0.72), col)
        txt(s, lab, Inches(0.7), y, Inches(0.9), Inches(0.72), size=15, color=WHITE,
            bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, Inches(1.7), y, Inches(10.9), Inches(0.72), LIGHT)
        txt(s, body, Inches(1.9), y, Inches(10.6), Inches(0.72), size=15.5,
            color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    cy = Emu(int(top) + len(steps) * int(Inches(0.92)) + int(Inches(0.1)))
    txt(s, conclusion, Inches(0.7), cy, Inches(11.9), Inches(0.9), size=18,
        color=GREEN, bold=True)
    notes(s, note); return s

def fit(s, path, cx, top, max_w, max_h):
    w, h = png_dim(path); ar = w / h
    W = max_w; H = Emu(int(W / ar))
    if H > max_h: H = max_h; W = Emu(int(H * ar))
    s.shapes.add_picture(path, Emu(int(cx - W / 2)), top, width=W, height=H)

def figure(title, img, howto, kicker="Results", note="", max_h=Inches(3.9)):
    s = prs.slides.add_slide(BLANK); heading(s, title, kicker)
    fit(s, img, SW // 2, Inches(1.8), Inches(11.4), max_h)
    rect(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(1.15), BOX)
    txt(s, "How to read it:  " + howto, Inches(0.95), Inches(6.15), Inches(11.4),
        Inches(1.0), size=14.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    notes(s, note); return s

def table(title, headers, rows, note_line, kicker="Results", hi=None, col_w=None,
          speaker=""):
    s = prs.slides.add_slide(BLANK); heading(s, title, kicker)
    nr, nc = len(rows) + 1, len(headers)
    tw = Inches(11.9); th = Inches(0.5 * nr)
    left = Emu(int((SW - tw) / 2)); top = Inches(1.95)
    t = s.shapes.add_table(nr, nc, left, top, tw, th).table
    if col_w:
        for j, w in enumerate(col_w): t.columns[j].width = Inches(w)
    for j, hh in enumerate(headers):
        c = t.cell(0, j); c.text = hh; c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
        rn = pr.runs[0]; rn.font.size = Pt(15); rn.font.bold = True; rn.font.color.rgb = WHITE
    for i, row in enumerate(rows, 1):
        h = (hi is not None and i - 1 == hi)
        for j, v in enumerate(row):
            c = t.cell(i, j); c.text = str(v); c.fill.solid()
            c.fill.fore_color.rgb = (RGBColor(0xE3,0xF1,0xE8) if h else (LIGHT if i % 2 else WHITE))
            pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            rn = pr.runs[0]; rn.font.size = Pt(14); rn.font.bold = h
            rn.font.color.rgb = GREEN if h else DARK
    rect(s, Inches(0.7), top + th + Inches(0.2), Inches(11.9), Inches(1.25), BOX)
    txt(s, "What it means:  " + note_line, Inches(0.95), top + th + Inches(0.3),
        Inches(11.4), Inches(1.1), size=14.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    notes(s, speaker); return s

def twocol(title, left_head, left_items, right_head, right_items, kicker=None, note=""):
    s = prs.slides.add_slide(BLANK); heading(s, title, kicker)
    for (x, head, items, col) in [(Inches(0.7), left_head, left_items, ACCENT),
                                  (Inches(6.95), right_head, right_items, NAVY)]:
        rect(s, x, Inches(1.85), Inches(5.65), Inches(0.55), col)
        txt(s, head, x, Inches(1.85), Inches(5.65), Inches(0.55), size=16, color=WHITE,
            bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb = s.shapes.add_textbox(x, Inches(2.55), Inches(5.65), Inches(4.4))
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            para(tf, it, size=15.5, first=(i == 0), bullet=True, after=8)
    notes(s, note); return s

# ================= TITLE =================
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
txt(s, "SERUM", Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.3), size=66, color=WHITE, bold=True)
txt(s, "Content-Aware Agentic Containment of Malware Under an Unobserved Payload",
    Inches(0.9), Inches(3.2), Inches(11.6), Inches(1.0), size=23, color=RGBColor(0xC7,0xD6,0xEC))
txt(s, "Semantic Epidemic Response under Unknown Malware", Inches(0.9), Inches(4.15),
    Inches(11.5), Inches(0.5), size=15, color=RGBColor(0x9F,0xB4,0xD4), italic=True)
txt(s, "Amritha S.", Inches(0.9), Inches(5.5), Inches(6), Inches(0.5), size=20, color=WHITE, bold=True)
txt(s, "A defensive research testbed at the intersection of cybersecurity, network "
    "science, and agentic AI", Inches(0.9), Inches(5.98), Inches(11), Inches(0.6),
    size=14, color=RGBColor(0x9F,0xB4,0xD4))
notes(s, "Welcome. SERUM is about containing a spreading computer worm when you don't "
      "know what vulnerability it's exploiting. I'll build the problem from scratch, "
      "give the method, the theory, the results, and be very honest about the limits.")

# ================= PART 1: THE PROBLEM =================
divider(1, "The Problem", "Why containing malware is hard when you can't see the attack")

explain("Malware spreads like an epidemic — but a picky one",
        "A self-propagating worm moves host to host, like a disease over a contact network.",
        [("Classic view: model it as an epidemic (SI/SIR) spreading over a graph of machines.",),
         ("The catch specific to malware: a worm exploiting vulnerability c can infect a "
          "neighbour ONLY IF that neighbour runs the vulnerable software (carries c).",),
         ("So infection is “vulnerability-gated”. A machine that doesn’t have the weakness "
          "is simply immune to THIS worm — the germ bounces off.",),
         ("Consequence: the worm does not travel the whole network. It travels only the "
          "sub-network of machines that share its target weakness.", 17, ACCENT)],
        kicker="The Problem",
        note="Start with the familiar epidemic-on-a-graph picture, then add the one twist that "
             "changes everything: infection only happens if the target machine is actually "
             "vulnerable. That gating is the whole basis of the project.")

worked("The picky germ, concretely",
       "Think of each machine as ‘wearing socks’ — the software it runs. The worm only "
       "infects one colour.",
       [("RED", "Machines running the vulnerable software (carry the target CVE) — CAN catch it.", GREEN),
        ("BLUE", "Machines running something else — the worm bounces off. Totally safe.", ACCENT),
        ("SPREAD", "The worm walks the network but only crosses into RED machines.", NAVY)],
       "So the ‘real’ network the worm uses = the RED sub-network, not the whole graph.",
       note="This is the kid-friendly version and it’s exactly right. Red socks = vulnerable "
            "software. The worm can only infect red machines. Keep this picture in mind — every "
            "later idea is about finding and defending the red machines.")

explain("The physical network is NOT the propagation graph",
        "The network a worm actually traverses is a payload-specific SUB-graph of the real network.",
        [("Physical topology: every machine and every link.",),
         ("Propagation graph: only the machines carrying the target CVE, and the links between them.",),
         ("These two graphs can look completely different — the vulnerable machines might be "
          "scattered in a corner, not spread across the hubs.",),
         ("A defender who confuses the two will defend the wrong thing.", 17, AMBER)],
        kicker="The Problem",
        note="The single most important conceptual point: don’t defend the physical graph, defend "
             "the propagation graph. The rest of the talk is about how to do that when you can’t "
             "directly see which machines are on it.")

explain("Why the usual defense wastes its budget",
        "Conventional defenses protect the most-connected machines (the ‘hubs’). Often that’s wasted.",
        [("You never have enough budget to protect every machine — you pick a few.",),
         ("Standard heuristics (degree, betweenness immunization) pick the network hubs.",),
         ("But a hub is irrelevant to THIS worm if it can’t run the exploit — it’s a blue-socks "
          "machine. Every band-aid spent on it is wasted.",),
         ("On real organisational networks, the vulnerable machines are frequently NOT the hubs — "
          "so ‘protect the hubs’ barely helps at all (we’ll see this holds empirically).", 17, ACCENT)],
        kicker="The Problem",
        note="This is the failure mode we beat. Structure-only defense is strong when the "
             "vulnerable set coincides with the hubs, but useless when it doesn’t — and real "
             "networks are often the latter.")

explain("The twist that makes it a real research problem",
        "The defender NEVER observes the payload. It only sees which machines are infected.",
        [("You don’t get told the worm’s target vulnerability — no signature, no payload capture.",),
         ("You observe: the set of currently-infected machines, and your own asset inventory "
          "(what software each machine runs).",),
         ("So you must DEFEND the vulnerable sub-network without being told which sub-network it is.",),
         ("That’s what elevates this from ‘SI on a subgraph’ to a genuine inference problem.", 17, NAVY)],
        kicker="The Problem",
        note="Emphasise the partial observability. If you knew the exploit, this would be easy — "
             "defend its subgraph. The difficulty and the contribution are about acting well "
             "WITHOUT that knowledge.")

explain("The clue hiding in plain sight",
        "Because spread is vulnerability-gated, every infected machine MUST carry the target CVE.",
        [("A machine can only have been infected by propagation if it was vulnerable to the worm.",),
         ("So each newly-infected machine is a hard CONSTRAINT: the true exploit lies somewhere in "
          "that machine’s software profile.",),
         ("Intersect those profiles across all infected machines → the set of exploits still "
          "consistent with what you’ve seen shrinks.",),
         ("SERUM turns this into online Bayesian inference — a POMDP whose hidden state is the "
          "exploit — and acts under the resulting belief.", 17, ACCENT)],
        kicker="The core idea",
        note="This is the ‘aha’. The infected set isn’t just damage — it’s evidence. Every "
             "infection narrows down the attack. This is the hook the whole method hangs on.")

# ================= PART 2: THESIS =================
divider(2, "Thesis & Contributions", "What we claim, and what we honestly do not")

s = prs.slides.add_slide(BLANK); bg(s, NAVY)
txt(s, "Thesis", Inches(0.8), Inches(0.7), Inches(6), Inches(0.5), size=18,
    color=RGBColor(0x9F,0xB4,0xD4), bold=True)
txt(s, "A defender that reasons about WHAT is spreading — the payload’s target "
    "vulnerability — contains a worm far more efficiently, with far less collateral "
    "disruption, than one that sees only network STRUCTURE.",
    Inches(0.8), Inches(1.5), Inches(11.8), Inches(2.6), size=27, color=WHITE, bold=True)
txt(s, "And this holds even when the payload is never directly observed — it is inferred "
    "from the shape of the outbreak.", Inches(0.8), Inches(4.5), Inches(11.8), Inches(1.2),
    size=20, color=RGBColor(0xC7,0xD6,0xEC), italic=True)
_n[0] += 1
notes(s, "The one-sentence thesis. Content beats structure, because malware can only traverse "
      "exploitable hosts — and we can pull it off without seeing the payload.")

twocol("Honest from the start: what is and isn’t new",
       "NOT novel (we claim none of this)",
       ["Vulnerability-gated spread = multitype bond percolation (known).",
        "Structure-only immunization baselines (degree, betweenness, …) are classical.",
        "POMDP defense, active hypothesis testing, LLM-as-prior — all established.",
        "The identifiability condition, once framed, is 1960s group-testing combinatorics."],
       "Our defensible contribution",
       ["The COUPLING no prior work occupies: online inference of the unobserved exploit "
        "from a vulnerability-gated cascade,",
        "with a checkable identifiability condition on OBSERVABLE host attributes,",
        "driving budgeted, content-aware containment,",
        "evaluated on REAL NVD/CVE data. (The cell CyGym 2025 leaves open.)"],
       kicker="Positioning",
       note="Disarm the novelty attack up front. Each ingredient is known; the combination and "
            "the real-data grounding are the contribution. Reviewers respect this honesty.")

explain("The four contributions (we lead with the finding)",
        "We foreground the empirical result and treat the inference and theory as its enablers.",
        [("C1 — Content-aware containment: on a REAL organisational topology, it substantially "
          "beats every structure-only baseline (−28.4%, p=1.7×10⁻⁷).", 16),
         ("C2 — An identifiability CHARACTERIZATION: identifiable iff infected profiles intersect "
          "in a singleton; plus a measured group-testing sample-complexity rate.", 16),
         ("C3 — A REAL-DATA evaluation (NVD/CVE) with paired statistics and an honest account of "
          "where content-awareness does and does not help.", 16),
         ("Robustness note — attacking the inference backfires; poisoning is absorbed (this is "
          "graceful degradation, not a superiority claim).", 16, GREY)],
        kicker="Overview",
        note="Four things. Notice we demote robustness to a note — it’s a null result (never worse "
             "than structure), not a win. Leading with the real-data result is deliberate.")

# ================= PART 3: FORMULATION =================
divider(3, "Problem Formulation", "The containment POMDP, precisely")

explain("The SERUM POMDP — the pieces",
        "A partially-observed control problem whose hidden state is the attacker’s exploit.",
        [("Hosts & profiles: graph G=(V,E); each host v carries X(v) ⊆ C, the set of CVEs it is "
          "exploitable by. CVE prevalence is heavy-tailed (a few common, many rare).",),
         ("Hidden state: the worm’s target c* ∈ C, drawn by the attacker. The defender never sees it.",),
         ("Vulnerable subgraph: carriers(c*) and the links among them — the graph the worm can use.",),
         ("Observation each step: the infected set (not c*) + the static inventory X.",),
         ("Actions (budget B/step): PATCH (immunize, keeps host online), ISOLATE (remove host, "
          "costs availability), SEGMENT (cut a link).", 17, DARK)],
        kicker="Method",
        note="Walk through each element slowly. The key line is ‘the defender never sees c*’ and "
             "‘every observation is the infected set’. That defines the inference problem.")

worked("How a single spread step works",
       "Infected u tries to infect susceptible neighbour w. Whether it succeeds is GATED:",
       [("CHECK", "Does w carry the target CVE c*?  (Is w red-socks?)", NAVY),
        ("NO", "w is immune to this worm. Nothing happens — ever.", ACCENT),
        ("YES", "w gets infected with probability β (transmissibility).", GREEN)],
       "Repeat over all infected→susceptible edges each step → the outbreak grows within the red subgraph.",
       note="Concretely: spread is a coin flip with probability β, but only for neighbours that are "
            "actually vulnerable. Blue neighbours never get infected. This gating is enforced in code "
            "and tested.")

explain("The objective — three things at once",
        "Containment is multi-objective; a good defender balances them, not just infection count.",
        [("Minimize final/peak INFECTION (how much of the fleet falls).",),
         ("Minimize AVAILABILITY loss (isolating machines takes them offline — collateral damage).",),
         ("Minimize TIME-TO-CONTAINMENT (stop the outbreak fast).",),
         ("This is why PATCH vs ISOLATE matters: patching immunizes without downtime, but you only "
          "dare patch precisely once you’re confident which machines are truly at risk.", 17, ACCENT)],
        kicker="Method",
        note="Real incident response isn’t just ‘fewest infections’ — taking critical machines "
             "offline is itself a cost. Our agent gets to keep more machines online because it "
             "knows WHICH ones to protect.")

# ================= PART 4: THEORY =================
divider(4, "Theory — Identifiability", "When can you actually figure out the exploit?")

explain("The question the theory answers",
        "Given only who’s infected, when is the hidden exploit uniquely recoverable?",
        [("If you can pin the exact CVE, you can defend its subgraph exactly.",),
         ("If two CVEs are indistinguishable from the infections seen, you must hedge across both.",),
         ("We want a CHECKABLE condition — one a defender can evaluate in advance from its inventory, "
          "before any outbreak.",),
         ("It turns out to be a simple set-containment condition on observable attributes.", 17, NAVY)],
        kicker="Theory",
        note="Frame the theory as answering a practical question: can I know the attack, or must I "
             "hedge? And crucially, can I tell in advance? Yes — from the asset inventory alone.")

worked("Worked example: narrowing to the exploit",
       "Watch the candidate set shrink as machines fall. supp = intersection of infected profiles.",
       [("Host A", "infected. Runs {CVE-1, CVE-2, CVE-5}.  Candidates = {1,2,5}.", NAVY),
        ("Host B", "infected. Runs {CVE-2, CVE-5, CVE-9}.  Candidates = {1,2,5}∩{2,5,9} = {2,5}.", NAVY),
        ("Host C", "infected. Runs {CVE-5, CVE-7}.  Candidates = {2,5}∩{5,7} = {5}.", GREEN)],
       "Intersection collapsed to a single CVE → the exploit is identified: it must be CVE-5.",
       note="This is the whole inference in one slide. Every infected host’s software profile is "
            "intersected in. When the intersection hits a single CVE, you’ve identified the attack. "
            "If it stalls at two or more, those are ‘confusers’.")

explain("The identifiability theorem",
        "c* is identifiable from a saturating outbreak over region R  ⟺  ∩(profiles in R) = {c*}.",
        [("‘Saturating’ = the outbreak has infected everything it can reach in the vulnerable subgraph.",),
         ("If the intersection is a single CVE, you know the attack exactly.",),
         ("If it’s two or more, those CVEs are ‘confusable’ — no amount of watching this outbreak "
          "separates them.",),
         ("Honest caveat: the belief’s surviving support EQUALS this intersection by construction, so "
          "this is a CHARACTERIZATION of the observation model, not a deep theorem. The 100% "
          "‘validation’ is a consistency check on the code, not an empirical prediction.", 16, AMBER)],
        kicker="Theory",
        note="State the theorem, then immediately own its modesty — this is exactly the honesty a "
             "committee rewards. The real value is the RATE (next), and knowing in advance which "
             "CVEs are confusable.")

explain("Confusability = the subset order",
        "One CVE hides behind another exactly when its victims are a subset of the other’s.",
        [("If carriers(c) ⊆ carriers(c′), then every c-victim also carries c′ — so a c-outbreak can "
          "never rule c′ out. c′ is a confuser of c.",),
         ("This gives a clean structure: order CVEs by carrier-set containment. A CVE is globally "
          "identifiable iff no other CVE’s carrier set is a superset of its own.",),
         ("Powerful consequence for DEFENSE: even if you can’t tell c from its confuser c′, defending "
          "c′’s (larger) subgraph OVER-COVERS the true victims — so you still protect them.", 16, GREEN)],
        kicker="Theory",
        note="This subset-order idea is why the method is robust: confusable exploits SHARE victims, "
             "so hedging over them still defends the truth. It’s also why an attacker can’t win by "
             "picking a confusable payload — more on that in the robustness section.")

explain("It’s group testing — realized by contagion",
        "The whole inference is an instance of online, graph-induced, adversarial group testing.",
        [("Classic group testing: identify a defective item by pooling and testing subsets.",),
         ("Here each infection is one ‘positive test’ whose profile-set is intersected into the "
          "running candidate set (Rényi separating systems; Kautz–Singleton cover-free families).",),
         ("The twist that IS ours: the ‘tests’ aren’t designed — they’re the fleet’s given software "
          "profiles — and the outcomes are produced by an adversarial, graph-constrained spreading "
          "process, observed online as the outbreak grows.",),
         ("So the number of infections needed inherits the group-testing log₂K bit-bound.", 16, ACCENT)],
        kicker="Theory",
        note="Placing it in group testing gives the theorem a rigorous home and preempts the "
             "information-theory reviewer. Our novelty is the online, contagion-realized, "
             "adversarial instantiation, not the separating condition itself.")

figure("Confusability structure (real NVD networks)",
       os.path.join(FIG, "confusability.png"),
       "left: most CVEs have ZERO confusers (identifiable). middle: the identifiable fraction "
       "decays as the CVE universe K grows (more CVEs → more subset-order dominators). right: a "
       "drawn slice of the confusability graph.",
       kicker="Theory",
       note="Three panels. Takeaway: on real data a majority of CVEs are individually identifiable, "
            "and identifiability gets harder as the vulnerability catalog grows — exactly as the "
            "subset-order theory predicts.")

figure("Sample complexity — how fast you identify the exploit",
       os.path.join(FIG, "sample_complexity.png"),
       "distribution of how many propagation infections it takes for the candidate set to collapse "
       "to one CVE. Median ≈ 5 infections ≈ 1.02·log₂K — essentially the information-theoretic "
       "group-testing minimum. Real profile correlation bends the rate TOWARD the optimum.",
       kicker="Theory",
       note="This is the non-trivial, testable content of the theory. log₂K is the theoretical "
            "minimum number of yes/no tests to find one item among K. We hit it — about 5 infected "
            "machines, ~1% of the fleet, identify the attack. That’s remarkably efficient.")

# ================= PART 5: METHOD =================
divider(5, "The Method", "The belief, and the content-aware policy")

worked("The belief update, step by step",
       "The agent keeps a probability over which CVE is the attack, and sharpens it each step.",
       [("PRIOR", "Start from CVE prevalence (or CVSS/LLM threat-intel for a cold start).", NAVY),
        ("OBSERVE", "New infections arrive. Each one must carry the true CVE.", ACCENT),
        ("UPDATE", "Down-weight (soft) or eliminate (hard) CVEs inconsistent with the new victims.", ACCENT),
        ("ACT", "Take the expectation of the defense score over the current belief.", GREEN)],
       "Early on the belief is broad → the agent hedges. As it sharpens → the agent commits.",
       note="The belief is just Bayesian filtering over the CVE. The important design choice is SOFT "
            "vs HARD, on the next slide. And note: the agent acts under UNCERTAINTY — it doesn’t wait "
            "to be sure.")

twocol("Soft vs hard likelihood — why soft wins",
       "HARD consistency (the theory)",
       ["A CVE absent from any infected host’s profile is ELIMINATED outright.",
        "Clean and exact — matches the identifiability theorem.",
        "Fragile: a single false-positive detection (a machine wrongly flagged infected) can "
        "wrongly eliminate the TRUE CVE. One bad observation breaks it."],
       "SOFT likelihood (what we deploy)",
       ["Each infected host DOWN-WEIGHTS rather than eliminates inconsistent CVEs.",
        "One false detection can’t exclude the truth — it just nudges the belief.",
        "Robust to noisy detection and to deliberate belief-poisoning.",
        "This is the version used in all deployment-facing experiments."],
       kicker="Method",
       note="Real detection is noisy. Hard elimination is brittle; soft down-weighting is robust. "
            "This single choice is what makes the agent survive noise and poisoning later.")

explain("The content-aware policy — the scoring rule",
        "Score each frontier host by its belief-weighted EXPOSED-VULNERABLE DEGREE.",
        [("For a host v, ask: how many of its neighbours are (a) still susceptible AND (b) exploitable "
          "by a CVE the belief thinks is live?",),
         ("That count — averaged over the belief — is v’s expected number of onward infections. Defend "
          "the hosts with the highest such score.",),
         ("It generalizes classic ‘degree’ immunization to the uncertain-payload regime, and reduces "
          "to plain degree when the belief is uniform / everything is vulnerable.",),
         ("Then: while uncertain, ISOLATE (a blunt but certain cut); once the belief’s support "
          "collapses, switch to PATCH (immunize precisely, keep hosts online).", 16, ACCENT)],
        kicker="Method",
        note="This is the heart of ‘content-aware’. It’s a payload-conditioned centrality: not ‘who "
             "is connected’ but ‘who can actually spread THIS worm next’. The isolate→patch switch is "
             "how it protects availability.")

# ================= PART 6: RESULTS =================
divider(6, "Results", "Synthetic, real-topology, and real-data — with honest scope")

explain("Experimental design — why it’s trustworthy",
        "Every policy faces the IDENTICAL outbreak (paired trials); numbers are statistically tested.",
        [("Paired design: same graph, same target CVE, same seeds, same infection coin-flips for "
          "every defender in a trial. This slashes variance and makes deltas real, not noise.",),
         ("Randomized target CVE per trial — no cherry-picked attack.",),
         ("Real vulnerabilities from the NVD/CVE database; per-CVE transmissibility from CVSS.",),
         ("Paired Wilcoxon tests + family-wise (Holm-Bonferroni) correction on the headline claims. "
          "We also report per-trial WIN RATES, not just means.", 16, NAVY)],
        kicker="Results",
        note="Sell the rigor before the numbers. Paired trials are the reason a 0.6-point mean "
             "difference can be highly significant — every defender fought the exact same outbreak.")

table("Headline — synthetic networks, real CVEs",
      ["Defender", "Infected ↓", "Availability ↑", "Contained@step ↓"],
      [["No defense", "36.6%", "100%", "14.2"],
       ["Best structural (degree / betweenness)", "3.4%", "95.5%", "4.5"],
       ["SERUM — content-aware (ours)", "1.8%", "98.5%", "3.0"],
       ["Oracle (is told the exploit)", "1.1%", "100%", "2.0"]],
      "Content-aware beats the best structural baseline on infection AND availability AND speed — "
      "without being told the attack — landing close to the oracle that IS told. Honest: the "
      "absolute gap here is small (1.8 vs 3.4 pts); the big, decisive gap shows up on real topology.",
      kicker="Results", hi=2, col_w=[6.3, 2.0, 2.1, 1.5],
      speaker="First headline. Note three axes, not one — we keep more machines online too. But be "
              "honest: on these synthetic networks the absolute numbers are small. The next slides "
              "are where it really matters.")

table("Flagship — a REAL organisational network",
      ["Defender", "Infected ↓", "Availability ↑"],
      [["No defense", "20.1%", "100%"],
       ["Degree immunization", "18.8%", "92.2%"],
       ["Betweenness (best structural)", "17.6%", "92.5%"],
       ["SERUM — content-aware (ours)", "11.7%", "97.5%"],
       ["Content-aware oracle (bound)", "9.1%", "100%"]],
      "SNAP email-Eu-core: 1004 REAL hosts, 42 REAL departments used as software zones, real NVD "
      "CVEs. Structure-only BARELY helps (17.6% vs 20.1% no-defense) — because the vulnerable "
      "departments are NOT the network hubs. Content-aware cuts infection to 11.7%: −28.4%, "
      "Wilcoxon p = 1.7×10⁻⁷. This is the result we lead with.",
      kicker="Results — flagship", hi=3, col_w=[7.7, 2.05, 2.05],
      speaker="THE result. On a real network with real community structure, ‘protect the hubs’ is "
              "almost worthless because the vulnerable machines aren’t the hubs — exactly the regime "
              "the thesis predicts. We cut infections by nearly a third, highly significant.")

figure("Flagship — infection over time (real topology)",
       os.path.join(FIG, "flagship_infection_curves.png"),
       "each curve = mean infected fraction over time for one defender. The structure-only curves sit "
       "right on top of no-defense; the content-aware curve is far below, hugging the oracle. The gap "
       "between them is the value of content-awareness.",
       kicker="Results — flagship",
       note="Visually: structure-only can’t separate from doing nothing here, while content-aware "
            "tracks the ‘cheating’ oracle. This picture IS the paper’s thesis.")

s = prs.slides.add_slide(BLANK); heading(s, "The lead result is NOT a fragile minority effect", "Results")
rect(s, Inches(0.7), Inches(2.0), Inches(5.2), Inches(3.6), BOX)
txt(s, "20 / 20", Inches(0.7), Inches(2.4), Inches(5.2), Inches(1.4), size=76,
    color=GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(s, "paired outbreaks WON on the real topology — vs betweenness AND vs the ensemble oracle "
    "(p = 8.8×10⁻⁵). 37/40 in a budget-8 replication.", Inches(0.9), Inches(3.9),
    Inches(4.8), Inches(1.6), size=16, color=DARK, align=PP_ALIGN.CENTER)
tb = s.shapes.add_textbox(Inches(6.3), Inches(2.0), Inches(6.4), Inches(4.6)); tf = tb.text_frame; tf.word_wrap=True
para(tf, "Why this slide exists:", size=18, bold=True, color=NAVY, first=True)
para(tf, "On the small synthetic margins, content-aware wins only a MINORITY of individual "
     "trials — the mean is carried by a few big wins. A fair critic would pounce on that.", size=16, bullet=True)
para(tf, "On the REAL topology it’s the opposite: it wins essentially EVERY paired outbreak, by "
     "5.9 points absolute. That’s why we lead with this, not the synthetic table.", size=16, bullet=True)
para(tf, "Honest boundary: on hub-aligned or near-universal-exploit regimes the edge is small or "
     "vanishes. We say so.", size=16, bullet=True, color=AMBER)
_n[0] += 0
notes(s, "Pre-empt the ‘tiny effect / you lose most trials’ attack head-on. It’s TRUE for the "
      "synthetic regime and we admit it — but the flagship, our actual headline, wins 20/20. "
      "Leading with the flagship is the honest, strong move.")

figure("Pareto dominance — infection vs availability",
       os.path.join(FIG, "pareto.png"),
       "each point = one defender at one budget. Down-and-right is better (less infection, more "
       "uptime). Content-aware sits in the bottom-right corner at EVERY budget — simultaneously "
       "lower-infection and higher-availability than every structural baseline.",
       kicker="Results",
       note="No trade-off is being hidden. At every budget, content-aware is better on BOTH axes. It "
            "achieves this by patching precisely once its belief sharpens, instead of isolating "
            "blindly.")

figure("The advantage scales with outbreak severity",
       os.path.join(RES, "prevalence_curve.png"),
       "relative infection reduction vs the best structural baseline, across exploit-prevalence "
       "bands. It grows with severity. Honest: in the low-prevalence bands content-aware wins a "
       "MINORITY of individual trials and one band isn’t significant; significance grows with severity.",
       kicker="Results",
       note="The worse the outbreak, the more content-awareness helps. But we report the uncomfortable "
            "part too: at low severity the per-trial win rate is below half. Honesty over hype.")

table("Head-to-head vs the CLOSEST prior systems",
      ["Defender (of this class)", "Infected ↓", "vs ours"],
      [["DAVA-style — data-aware but exploit-blind", "1.70%", "+43.8%  ·  p=2.8×10⁻⁴"],
       ["CyGym-style — static prior, no online update", "1.15%", "+16.6%  ·  p=1.1×10⁻²"],
       ["SERUM — content-aware (ours)", "0.95%", "—"]],
      "We beat BOTH nearest systems. DAVA even underperforms plain degree — vaccinating "
      "exposed-but-non-exploitable hosts wastes budget, i.e. the thesis. Honest scope: these are "
      "faithful reimplementations of each system’s STANCE, not the full originals — a fair caveat "
      "we state, with a released-code port as the follow-up.",
      kicker="Results — vs prior work", hi=2, col_w=[7.0, 2.0, 2.9],
      speaker="This answers ‘did you compare to the real competitors?’. Yes. Content beats "
              "data-aware (DAVA) and static-prior (CyGym). We’re upfront that these are "
              "reimplementations of their epistemic stance, not the full systems.")

explain("When does the ONLINE inference actually matter?",
        "Honest finding: content-AWARENESS is the driver; the online UPDATE is a smaller refinement.",
        [("Under a GOOD prior, online updating beats a frozen (static) prior by only +0.19pp — "
          "consistent with our ablation that freezing the belief costs almost nothing.",),
         ("It earns its keep when the prior is WRONG: under a deliberately misleading prior, the "
          "edge roughly DOUBLES to +0.44pp (p=1.8×10⁻²).",),
         ("So online inference is a refinement that matters most under bad threat intel — NOT the "
          "main source of the win. We say this plainly rather than overselling ‘online inference’.", 16, AMBER)],
        kicker="Analysis",
        note="A subtle, honest point. Our own experiments show the online update is a small effect "
             "under a good prior. The real win is defending the vulnerable subgraph at all "
             "(content-awareness). We reframed the paper to say this.")

# ================= PART 7: ROBUSTNESS =================
divider(7, "Robustness", "Can an attacker defeat the defender by attacking its inference?")

explain("Attacking the inference BACKFIRES",
        "An attacker who picks a hard-to-identify (confusable) payload does not help itself.",
        [("Intuition from the subset-order theory: a confusable exploit’s victims are a SUBSET of "
          "its confuser’s victims.",),
         ("So even a hedged belief that can’t tell them apart still DEFENDS the true victims — they’re "
          "inside the set it’s protecting.",),
         ("Empirically the content-aware edge is actually LARGER under an evasive attacker (+22.4%) "
          "than a random one (+17.7%).",),
         ("Inference-evasion is not a winning strategy. This is a structural guarantee, not luck.", 16, GREEN)],
        kicker="Robustness",
        note="The attacker’s natural move — hide behind a common CVE — doesn’t work, because hiding "
             "means sharing victims, and we defend shared victims anyway. Elegant and it’s proven, "
             "not just measured.")

explain("Belief poisoning, and the honest limit",
        "A stronger attacker plants FAKE infections to mislead the belief. Here’s what actually happens.",
        [("The fix: a RobustAgent audits its belief against the real spread and hedges toward "
          "structure-only defense when they stop matching — poisoning is one-shot, the real worm "
          "keeps revealing the truth.",),
         ("Against a white-box, audit-AWARE adaptive poisoner, it holds up to a poisoning budget of "
          "6% of the fleet; only at an extreme 10% does the attacker gain a small edge (which fails "
          "multiple-comparison correction).",),
         ("HONEST framing: this is GRACEFUL DEGRADATION, not a robustness win. Under poisoning the "
          "content-aware advantage evaporates and the agent falls back to — never below — "
          "structure-only. It’s a safety net, not superiority under attack.", 15, AMBER)],
        kicker="Robustness",
        note="This is where we’re most careful not to overclaim. The robust agent guarantees you’re "
             "never WORSE than a dumb heuristic under poisoning. That’s worth stating — but it’s a "
             "null result, not a victory, and we label it as such.")

twocol("Two more robustness stresses",
       "Imperfect detection sensors",
       ["Model missed detections + persistent false alarms.",
        "Content-aware degrades GRACEFULLY: positive edge at 9 of 10 noise points.",
        "False alarms (which poison the belief) are the sensitive channel.",
        "Not a new significant win — a graceful-degradation result."],
       "Is the win a manufactured artifact?",
       ["Worry: the ‘homophily’ knob manufactures the favorable regime.",
        "We sweep it 0→0.8. Surprise: the edge is significant even at 0 (no monoculture).",
        "So it’s NOT a knob artifact — homophily only controls spatial clustering.",
        "Honest: this rebuts the KNOB worry, not the semi-synthetic-data limitation."],
       kicker="Robustness / threats to validity",
       note="Two stress tests. Left: noisy sensors — degrades gracefully. Right: we tried to break "
            "our OWN result by turning off the monoculture knob, and the advantage survived. That "
            "rebuttal makes the result stronger.")

# ================= PART 8: BREADTH =================
divider(8, "Breadth", "The framework generalizes")

twocol("Application & learned policy",
       "IoT botnet (Mirai-style)",
       ["Device-type firmware monoculture; payload targets default telnet credentials.",
        "Host value = uplink Mbps, so ‘blast radius’ = the DDoS capacity a botnet conscripts.",
        "Content-aware cuts DDoS blast −8.71pp vs degree (20/20 wins), higher availability.",
        "Absolute infection stays high — tight-budget Mirai does form botnets; read the deltas."],
       "Learned policy (validation, not headline)",
       ["A cross-entropy-method policy over belief-augmented features…",
        "…matches the hand-designed agent (1.09% vs 1.04%) and beats degree/betweenness.",
        "It independently learns to weight the belief features most — validating the design.",
        "This is learning UNDER exploit uncertainty, not vanilla GNN-RL immunization."],
       kicker="Breadth",
       note="Two breadth results. IoT shows the framework instantiates on a real threat model. The "
            "learned policy shows a from-scratch learner rediscovers our hand-designed signal — "
            "evidence the design is right, not a separate contribution.")

# ================= PART 9: POSITIONING & HONESTY =================
divider(9, "Positioning & Honesty", "Nearest work, and the limits we own")

table("The three nearest systems — cite, then differentiate",
      ["Prior work", "What it shares with SERUM", "The gap SERUM fills"],
      [["CyGym (2025)", "vuln-gated spread + a cost model",
        "STATIC zero-day prior; no online belief update"],
       ["SCENARIOID (KDD 2023)", "infer which mechanism produced a cascade",
        "OFFLINE, feature-engineered; no identifiability guarantee"],
       ["Hoffmann et al. (ICML 2020)", "identifiability from cascades",
        "condition on LATENT edges; batch (an arguably harder setting)"]],
      "We cite all three prominently and state each delta in the same breath — including, honestly, "
      "that our observable-attribute setting is in some ways EASIER than Hoffmann’s latent one, not "
      "strictly stronger. Reviewers pattern-match to these three; we get ahead of it.",
      kicker="Related work", col_w=[2.9, 4.5, 4.5],
      speaker="These are the papers a committee will say ‘isn’t this just…’. We name them first and "
              "state exactly what’s different — and where we’re actually the easier problem. Honesty "
              "buys credibility.")

explain("Honest limitations (several are load-bearing)",
        "We state these plainly; a paper is stronger for owning them.",
        [("L1 — Not yet validated on real HOST-LEVEL data: measured per-host inventories on a real "
          "segmented network are proprietary scan data. The single most important gap.", 15),
         ("L2 — Online inference is a REFINEMENT, not the driver (our own ablation shows this).", 15),
         ("L3 — The host↔CVE mapping is MODELED (semi-synthetic), not measured.", 15),
         ("L4 — Detection noise: addressed (graceful degradation). L5 — Adaptive adversary: "
          "addressed, bounded. L6 — Multiplicity: corrected, but forking paths remain. L7 — Abstract "
          "model: no packet-level malware, scanning, or C2 layer.", 15)],
        kicker="Honesty",
        note="Don’t bury limitations — lead with them. L1 is the big one. Note we cross-reference "
             "which are addressed (L4, L5, L6) and which are genuine open scope (L1, L3, L7).")

explain("L1 — ready to close in ONE command",
        "We can’t close L1 without proprietary data — but we made the project ready the instant it exists.",
        [("We will NOT fabricate a real inventory, and another synthetic model wouldn’t satisfy L1.",),
         ("Instead we ship a TESTED import pipeline: serum/data/real_inventory.py + "
          "scripts/validate_real_inventory.py.",),
         ("Feed it a real scan (measured host,CVE findings) + a topology edge list, and it runs the "
          "full evaluation — taking each host’s vulnerabilities VERBATIM from the scan (measured, "
          "not modeled). Every downstream component works unchanged.",),
         ("So L1 is gated on DATA ACCESS, not on method or engineering. One command away.", 16, GREEN)],
        kicker="Honesty",
        note="Turn the biggest weakness into a readiness story. We didn’t fake the result; we built "
             "the exact, tested pipeline that produces it the moment someone shares a Nessus/Tenable "
             "export. Great answer to ‘so you have no real result’.")

explain("What holds up under grilling",
        "This was hardened through five rounds of hostile self-review.",
        [("Pareto dominance: at least as good on infection, strictly better on availability & speed.",),
         ("Real-topology flagship: −28.4%, p=1.7×10⁻⁷, wins 20/20 paired outbreaks.",),
         ("Beats the closest prior systems (CyGym-style, DAVA-style) head-to-head.",),
         ("Every headline number is backed by a committed result file and guarded by a test; the "
          "experiments are seeded and reproduce bit-for-bit; the paper compiles clean; 141 tests green.", 16, GREEN)],
        kicker="Summary",
        note="Close the technical arc with what survived scrutiny. The reproducibility and "
             "test-guarding matter to reviewers — nothing here is a number we can’t regenerate on "
             "demand.")

# ================= PART 10: CLOSE =================
divider(10, "Outlook", "")

explain("Future work",
        "The method is frozen and honest; the frontier is data and adversaries.",
        [("Validate on a real host-level enterprise inventory (L1) — the highest-value next step.",),
         ("A jointly-adaptive adversary that best-responds on payload + timing + placement together.",),
         ("End-to-end containment for polymorphic (multi-exploit) worms.",),
         ("A learned GNN policy and an LLM tool-use agent that reads structured threat intel.",),
         ("An emulation bridge (packet-level) from the abstract epidemiological model.",)],
        kicker="Outlook",
        note="The most important item is L1 — real data. Everything else is depth. We’re transparent "
             "that the exciting open work is external (data) and adversarial, not more tuning.")

s = prs.slides.add_slide(BLANK); bg(s, NAVY)
txt(s, "Defend what’s spreading —", Inches(0.9), Inches(2.2), Inches(11.6), Inches(1.0),
    size=42, color=WHITE, bold=True)
txt(s, "not just who’s connected.", Inches(0.9), Inches(3.25), Inches(11.6), Inches(1.0),
    size=42, color=RGBColor(0xC7,0xD6,0xEC), bold=True)
txt(s, "Infer the unseen exploit from the shape of the outbreak, and defend the machines "
    "that can actually catch it.", Inches(0.9), Inches(4.55), Inches(11.6), Inches(1.0),
    size=18, color=RGBColor(0x9F,0xB4,0xD4), italic=True)
txt(s, "SERUM — Semantic Epidemic Response under Unknown Malware   ·   Amritha S.   ·   "
    "Defensive research; no weaponizable attack code.", Inches(0.9), Inches(6.2),
    Inches(11.6), Inches(0.6), size=13, color=RGBColor(0x7F,0x94,0xB4))
_n[0] += 1
notes(s, "Close on the one-line takeaway. Thank the audience; invite the hardest question — "
      "especially about L1, which we’re ready for.")

prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
